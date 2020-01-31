from pptx import Presentation
from pptx.util import Cm

from PIL import Image

prs = Presentation('Jumbo 3.pptx')

# def delete_slides(presentation, index):
# 	xml_slides = presentation.slides._sldIdLst  
# 	slides = list(xml_slides)
# 	print(len(xml_slides))
# 	xml_slides.remove(slides[index])

slide = prs.slides[3]

for shape in slide.shapes:
	if hasattr(shape, "image"):
		img = shape
		if img.width > 2000000 : 
			image = "IMG_0932.jpg"
			imgPic = img._pic
			imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
			imgPart = slide.part.related_parts[imgRID]

			with open(image, 'rb') as f:
			    rImgBlob = f.read()
			rImgWidth, rImgHeight = Image.open(image).size
			rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight) # change from Px

			# replace
			imgPart._blob = rImgBlob
			
			widthScale = float(rImgWidth) / img.width
			heightScale = float(rImgHeight) / img.height
			maxScale = max(widthScale, heightScale)
			scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
			# center the image if it's different size to the original
			scaledImgLeft = int(img.left + (img.width - scaledImgWidth)/2)
			scaledImgTop = int(img.top + (img.height - scaledImgHeight)/2)
			# now update
			img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
			break

# for slide in prs.slides:
	# print(slide.slide_id)
	# for shape in slide.shapes:

		# if hasattr(shape, "text"):
		# 	# shape.text_frame.paragraphs[0].runs[0].text
		# 	paras = shape.text_frame.paragraphs
		# 	for i in range(len(paras)):
		# 		runs = paras[i].runs
				
		# 		if len(runs) > 1:
		# 			runs[0].text = paras[i].text
		# 		# print(i,paras[i].text)
		# 			for j in range(1,len(runs)):
		# 				runs[j].text = ""
		# 				# if i == 5 and j == 0:
		# 				# 	shape.text_frame.paragraphs[i].runs[j].text += "Added"
		# 				# print(i,j,runs[j].text)
		# 		for j in range(len(runs)):
		# 			print(i,j,runs[j].text)

# 	print()

# delete_slides(prs,0)

# def move_slide(p,old_index, new_index):
# 	xml_slides = p.slides._sldIdLst  # pylint: disable=W0212
# 	slides = list(xml_slides)
# 	xml_slides.remove(slides[old_index])
# 	xml_slides.insert(new_index, slides[old_index])


# move_slide(prs,0,2)
# move_slide(prs,1,-1)
# move_slide(prs,2,-1)


prs.save("edited.pptx")





