from pptx import Presentation
import copy
import six
import re
from PIL import Image
from pptx.util import Cm


class Document(object):
	"""docstring for Document"""
	def __init__(self, docname):
		super(Document, self).__init__()
		
		self.docname = docname
		self.doc = Presentation(self.docname)

	def delete(self,slideno):
		xml_slides = self.doc.slides._sldIdLst  
		slides = list(xml_slides)
		xml_slides.remove(slides[slideno])

	def deletebySlideID(self,slide_id):

		slide = self.doc.slides.get(slide_id)
		index = self.doc.slides.index(slide)
		self.delete(index)

	def deletebyTitle(self,title):
		slidestoDelete = []
		i = 0
		for slide in self.doc.slides:
			for shape in slide.shapes:
				if hasattr(shape, "text"):
					paras = shape.text_frame.paragraphs
					if paras[0].text == title:
						slidestoDelete.append(i)
						break
			i+=1

		print("Slides to Delete:",slidestoDelete)

		for j in range(len(slidestoDelete)):
			# re-adjusting index to accomodate reindexing of slide list
			newindex = slidestoDelete[j] - j
			self.delete(newindex)

	def blueprint(self):

		blueprint = {}
		slideid = 0
		for slide in self.doc.slides:
			slideid+=1
			slide_dict = {}
			seenOne = False
			desc_count = 0
			for shape in slide.shapes:
				if hasattr(shape, "text"):
					paras = shape.text_frame.paragraphs
					slide_dict["type"] = "Details"
					slide_dict["elements"] = []
					for i in range(len(paras)):
						if seenOne == False and len(paras[i].text) > 1 and len(paras[i].text) < 50: 
							# this detects the slide title
							seenOne = True
							slide_dict["title"] = paras[i].text
							slide_dict["id"] = str(slide.slide_id)
							element_dict = {}
							element_dict["key"] = "Title"
							element_dict["value"] = paras[i].text
							element_dict["pid"] = i
							slide_dict["elements"].append(element_dict)

						elif " Image" in paras[i].text:
							# this detected the slide type
							slide_dict["type"] = paras[i].text
							element_dict = {}
							element_dict["key"] = "Slide Type"
							element_dict["value"] = paras[i].text
							element_dict["pid"] = i
							slide_dict["elements"].append(element_dict)

						elif ":" in paras[i].text.lower() and len(paras[i].text) < 50:
							# this is detected as a key value entity
							element_dict = {}
							element_dict["key"] = paras[i].text.split(":")[0]
							element_dict["value"] = paras[i].text.split(":")[1]
							element_dict["pid"] = i
							slide_dict["elements"].append(element_dict)

						elif len(paras[i].text) > 50:
							# this is detected as a paragraph
							desc_count += 1
							element_dict = {}
							element_dict["key"] = "Description "+str(desc_count)
							element_dict["value"] = paras[i].text
							element_dict["pid"] = i
							slide_dict["elements"].append(element_dict)

			if "id" not in slide_dict:
				slide_dict["id"] = str(slide.slide_id)
				slide_dict["type"] = "Single Slide"
				slide_dict["elements"] = []
				slide_dict["title"] = "Intro Slide "+str(slideid)
				blueprint["Intro Slide "+str(slideid)] = [slide_dict]

			else:
				if slide_dict["title"] not in blueprint:
					blueprint[slide_dict["title"]] = []

				blueprint[slide_dict["title"]].append(slide_dict)


		return blueprint

	def stripPPT(self,selectedstuff):
		'''
		Function to remove all unnecessary sections from ppt
		and keep only the selected slides, and content,
		do so by sending a list of slide IDs and removing all
		slides not in that list.
		Similarly by sending all the content ids, and removing
		everything not in that list.
		'''

		slide_order = list(map(int,selectedstuff["slide_order"].split("_")))

		slide_ids_to_keep = []
		text_ids_to_keep = {}

		blueprint = self.blueprint()
		
		for key in selectedstuff.keys():

			if key[0] == "A":
				# keys that start with A_ represent sections to keep
				# and their values contain slide_IDs
				slide_ids_to_keep.append(int(selectedstuff[key]))

			elif key[0] == "B":
				# keys that start with B_ represent paragraphs to keep
				# and their values contain the actual values of the paragraph
				# format: B_slideid_textid_textkey
				value_split = key.split("_")
				slide_id = value_split[1]
				text_id = value_split[2]
				text_key = value_split[3]

				# add slideid key to dict
				if slide_id not in text_ids_to_keep:
					text_ids_to_keep[slide_id] = {}

				# create a heirarchieal datastructure to represent the text
				# to keep and what value to store
				text_ids_to_keep[slide_id][text_id] = [text_key,selectedstuff[key]]

		# remove unselected text from slides 
		for slide_id in text_ids_to_keep:
			for shape in self.doc.slides.get(int(slide_id)).shapes:
				if hasattr(shape, "text"):
					paras = shape.text_frame.paragraphs
					
					i = 0
					# i tracks the paragraph index in the textbox
					# we keep track of which paragraph index are selected
					for para in paras:
						if ((str(i) not in text_ids_to_keep[slide_id]) and 
							(para.text != "")):
							
							# actual delete paragraph ops
							p = para._p
							p.getparent().remove(p)

						elif (str(i) in text_ids_to_keep[slide_id]):
							
							value = text_ids_to_keep[slide_id][str(i)][1]

							if ":" in para.text:
								runs = para.runs
								# we assume there is only single run
								if len(runs) > 1:
									
									runs[0].text = para.text
									runs[0].text = runs[0].text.replace("_x000B_","")
									
									for j in range(1,len(runs)):
										runs[j].text = ""

								index = para.runs[0].text.index(":")
								paras[i].runs[0].text = paras[i].runs[0].text[:index+1] + value
								
							else:
								# runs = para.runs
								# if len(runs) > 1:
								# 	runs[0].text = para.text
								# 	for j in range(1,len(runs)):
								# 		runs[j].text = ""
								
								# paras[i].runs[0].text = value
								self.replace_paragraph_retaining_formatting(paras[i],value)
						i+=1



		# remove slides that are not selected
		for section in blueprint.keys():
			for slide in blueprint[section]:
				if int(slide["id"]) not in slide_ids_to_keep:
					self.deletebySlideID(int(slide["id"]))

		i = 0
		for slide_id in slide_order:

			slide = self.doc.slides.get(slide_id)
			index = self.doc.slides.index(slide)

			self.move_slide(index, i)
			i+=1

	def move_slide(self, old_index, new_index):
		xml_slides = self.doc.slides._sldIdLst  # pylint: disable=W0212
		slides = list(xml_slides)
		xml_slides.remove(slides[old_index])
		xml_slides.insert(new_index, slides[old_index])
	
	def replace_paragraph_retaining_formatting(self,paragraph, new_text):
	    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
	    # remove all but the first run
	    for idx, run in enumerate(paragraph.runs):
	        if idx == 0:
	            continue
	        p.remove(run._r)
	    if len(paragraph.runs) > 0:
		    paragraph.runs[0].text = new_text

	def editTitle(self,slide,value):

		for shape in slide.shapes:
			if hasattr(shape, "text"):
				paras = shape.text_frame.paragraphs
				for i in range(len(paras)):
					if len(paras[i].runs) > 0:
						self.replace_paragraph_retaining_formatting(paras[i],value)
						return True

		return False

	def editKeyValue(self,slide,key,value):

		for shape in slide.shapes:
			if hasattr(shape, "text"):
				paras = shape.text_frame.paragraphs
				for i in range(len(paras)):
					r = paras[i].runs
					if ((":" in paras[i].text) and 
						(key.lower() in paras[i].text.lower())):

						# index = paras[i].text.index(":")
						# we assume there is only single run
						# paras[i].runs[0].text = paras[i].runs[0].text[:index+1] + value
						
						self.replace_paragraph_retaining_formatting(paras[i],key+value)
						
						return True

		return False

	def editValue(self,slide,oldvalue,newvalue):
		for shape in slide.shapes:
			if hasattr(shape, "text"):
				paras = shape.text_frame.paragraphs
				for i in range(len(paras)):

					if paras[i].text == oldvalue:
						self.replace_paragraph_retaining_formatting(paras[i],newvalue)

					return True

		return False

	def editParagraph(self,slide,content):
		for shape in slide.shapes:
			if hasattr(shape, "text"):
				paras = shape.text_frame.paragraphs
				for i in range(len(paras)):
					if len(paras[i].text) > 50 :
						# qualifies as paragraph
						
						self.replace_paragraph_retaining_formatting(paras[i],content)

						return True
		return False

	def save(self,name):
		
		self.doc.save(name)

	def replace_image(self,slide,imgpath):
		for shape in slide.shapes:
			if hasattr(shape, "image"):
				img = shape
				if img.width > 2000000 : 
					imgPic = img._pic
					imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
					imgPart = slide.part.related_parts[imgRID]

					with open(imgpath, 'rb') as f:
					    rImgBlob = f.read()
					rImgWidth, rImgHeight = Image.open(imgpath).size
					rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight) # change from Px

					# replace
					imgPart._blob = rImgBlob
					
					widthScale = float(rImgWidth) / img.width
					heightScale = float(rImgHeight) / img.height
					maxScale = max(widthScale, heightScale)
					scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
					# center the image if it's different size to the original
					scaledImgLeft = int(img.left + (img.width - scaledImgWidth)/2)
					scaledImgTop = int(img.top + (img.height - scaledImgHeight)/2)
					# now update
					img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
					break


	def create_slide(self,details,imgpath):

		newslide = None

		if details["stone_imagetype"] == "":
			newslide = self.duplicate_slide(3)
		else:
			newslide = self.duplicate_slide(6)

		for key in details:

			if key == "stone_name":
				self.editTitle(newslide, details[key])

			elif key == "stone_description" and details["stone_imagetype"] == '':
				self.editParagraph(newslide,details[key])

			elif key == "stone_imagetype" and details["stone_imagetype"] != '' :
				self.editValue(newslide,"Bookmatch Image",details[key])

			elif key == "stone_price":
				self.editKeyValue(newslide,"Price:",details[key])

			elif key == "stone_discount":
				self.editKeyValue(newslide,"Discount:",details[key])

			elif key == "stone_discounted_price":
				self.editKeyValue(newslide,"Discounted Price:",details[key])

			elif key == "stone_size":
				self.editKeyValue(newslide,"Size:",details[key])

			elif key == "stone_quantity":
				self.editKeyValue(newslide,"Quantity:",details[key])

			elif key == "stone_processing":
				self.editKeyValue(newslide,"Processing:",details[key])

			elif key == "stone_thickness":
				self.editKeyValue(newslide,"Thickness:",details[key])


		self.replace_image(newslide,imgpath)

		self.save("Jumbo 3.pptx")


	def duplicate_slide(self,index):

		template = self.doc.slides[index]
		try:
			blank_slide_layout = self.doc.slide_layouts[4]
		except:
			blank_slide_layout = self.doc.slide_layouts[len(self.doc.slide_layouts)-1]

		copied_slide = self.doc.slides.add_slide(blank_slide_layout)

		for shp in template.shapes:
			el = shp.element
			newel = copy.deepcopy(el)
			copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

		for _, value in six.iteritems(template.part.rels):
			# Make sure we don't copy a notesSlide relation as that won't exist
			if "notesSlide" not in value.reltype:
				copied_slide.part.rels.add_relationship(value.reltype,
												value._target,
												value.rId)

		return copied_slide


# TODO: 
# Adding tags to ppt elements, for template creation
# PPT Versioning for various clients, and history of what was sent when and to whom
# 



